"""Tool implementations the agent can call.

Every capability the assistant has (web search, scraping, browser control,
PDF reading, document memory, desktop control, file creation) lives here as a
plain function. The agent loop in agent.py decides when to call them — there
is no keyword routing anywhere.
"""

import io
import json
import os
import re
import subprocess
from html import unescape
from pathlib import Path
from urllib.parse import quote_plus, unquote
from uuid import uuid4

import pdfplumber
import requests
from PIL import Image
from selenium.webdriver.common.by import By
from sqlalchemy.orm import Session

from models import DocumentChunk, ImageAttachment
from scraper import create_driver, extract_mcq_questions, fetch_page
from vector_store import retrieve_relevant_chunks, split_text

UPLOAD_DIR = Path(__file__).resolve().parent / "uploads"
GENERATED_DIR = Path(__file__).resolve().parent / "generated"

ACTIVE_BROWSER_SESSIONS = {}
ACTIVE_APP_SESSIONS = {}

DEFAULT_APP_REGISTRY = {
    "microsoft word": {"command": "winword", "exe": "WINWORD.EXE"},
    "word": {"command": "winword", "exe": "WINWORD.EXE"},
    "microsoft excel": {"command": "excel", "exe": "EXCEL.EXE"},
    "excel": {"command": "excel", "exe": "EXCEL.EXE"},
    "powerpoint": {"command": "powerpnt", "exe": "POWERPNT.EXE"},
    "power point": {"command": "powerpnt", "exe": "POWERPNT.EXE"},
    "chrome": {"command": "chrome"},
    "edge": {"command": "msedge"},
    "notepad": {"command": "notepad"},
    "calculator": {"command": "calc"},
    "calc": {"command": "calc"},
    "camera": {"uri": "microsoft.windows.camera:"},
    "photos": {"uri": "ms-photos:"},
    "settings": {"uri": "ms-settings:"},
    "mail": {"uri": "outlookmail:"},
    "calendar": {"uri": "outlookcal:"},
    "store": {"uri": "ms-windows-store:"},
    "vscode": {"command": "code"},
    "vs code": {"command": "code"},
}

# How much of a tool result the model sees. Full text is always saved to the
# chat's document memory, so nothing is lost by truncating here.
TOOL_RESULT_CHAR_LIMIT = 6000


def app_registry():
    raw = os.getenv("JARVIS_APP_REGISTRY")
    registry = dict(DEFAULT_APP_REGISTRY)
    if raw:
        try:
            loaded = json.loads(raw)
            if isinstance(loaded, dict):
                registry.update(loaded)
        except json.JSONDecodeError:
            pass
    return registry


def clip(text: str, limit: int = TOOL_RESULT_CHAR_LIMIT):
    text = text.strip()
    if len(text) <= limit:
        return text
    return text[:limit].rstrip() + "\n...[truncated — full text is saved in chat memory, use search_documents to find specific parts]"


# ---------------------------------------------------------------------------
# Document memory (uploaded files, scraped pages)
# ---------------------------------------------------------------------------

def extract_file_text(filename: str, file_bytes: bytes):
    lower_name = filename.lower()

    if lower_name.endswith(".pdf"):
        content = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    content.append(text)
        return "\n".join(content)

    if lower_name.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff")):
        import pytesseract

        image = Image.open(io.BytesIO(file_bytes))
        return pytesseract.image_to_string(image)

    if lower_name.endswith(".docx"):
        from docx import Document

        document = Document(io.BytesIO(file_bytes))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)

    if lower_name.endswith((".xlsx", ".xlsm")):
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
        rows = []
        for sheet in workbook.worksheets:
            rows.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value) for value in row if value is not None]
                if values:
                    rows.append(" | ".join(values))
        return "\n".join(rows)

    if lower_name.endswith((".doc", ".xls")):
        raise ValueError(
            "Old .doc/.xls files are not supported yet. Save as .docx or .xlsx and upload again."
        )

    return file_bytes.decode("utf-8", errors="replace")


def is_image_filename(filename: str):
    return filename.lower().endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"))


def save_image_attachment(db: Session, chat_id: str, filename: str, file_bytes: bytes):
    suffix = Path(filename).suffix.lower() or ".jpg"
    image_dir = UPLOAD_DIR / chat_id
    image_dir.mkdir(parents=True, exist_ok=True)
    image_path = image_dir / f"{uuid4()}{suffix}"
    image_path.write_bytes(file_bytes)

    db.add(ImageAttachment(
        id=str(uuid4()),
        chat_id=chat_id,
        filename=filename,
        path=str(image_path),
    ))
    return image_path


def save_document_text(db: Session, chat_id: str, filename: str, content: str):
    chunks = split_text(content.strip())
    if not chunks:
        return 0

    db.add_all([
        DocumentChunk(id=str(uuid4()), chat_id=chat_id, filename=filename, content=chunk)
        for chunk in chunks
    ])
    return len(chunks)


def list_chat_documents(db: Session, chat_id: str):
    rows = db.query(DocumentChunk.filename).filter(DocumentChunk.chat_id == chat_id).distinct().all()
    return [row[0] for row in rows]


def search_chat_documents(db: Session, chat_id: str, query: str, limit: int = 5):
    chunks = [
        f"Source: {row.filename}\n{row.content}"
        for row in db.query(DocumentChunk).filter(DocumentChunk.chat_id == chat_id).all()
    ]
    if not chunks:
        return "This chat has no saved documents yet. Ask the user to upload a file or use fetch_webpage/read_browser_page to add one."
    relevant = retrieve_relevant_chunks(chunks, query, limit=limit)
    return clip("\n\n---\n\n".join(relevant))


def build_context(db: Session, chat_id: str, query: str):
    chunks = [
        f"Source file: {row.filename}\nExtracted text:\n{row.content}"
        for row in db.query(DocumentChunk).filter(DocumentChunk.chat_id == chat_id).all()
    ]
    return "\n\n".join(retrieve_relevant_chunks(chunks, query))


def get_latest_image_payloads(db: Session, chat_id: str, limit=1):
    import base64

    rows = db.query(ImageAttachment).filter(ImageAttachment.chat_id == chat_id).all()
    payloads = []
    for row in rows[-limit:]:
        image_path = Path(row.path)
        if image_path.exists():
            payloads.append(base64.b64encode(image_path.read_bytes()).decode("utf-8"))
    return payloads


# ---------------------------------------------------------------------------
# Web: search, headless fetch, PDF from URL
# ---------------------------------------------------------------------------

def web_search(query: str, max_results: int = 6):
    """Search DuckDuckGo and return real results (title, url, snippet)."""
    response = requests.post(
        "https://html.duckduckgo.com/html/",
        data={"q": query},
        headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"},
        timeout=20,
    )
    response.raise_for_status()
    html = response.text

    titles = re.findall(r'<a[^>]+class="result__a"[^>]+href="([^"]+)"[^>]*>(.*?)</a>', html, re.S)
    snippets = re.findall(r'<a[^>]+class="result__snippet"[^>]*>(.*?)</a>', html, re.S)

    results = []
    for index, (href, raw_title) in enumerate(titles[:max_results]):
        url = href
        redirect = re.search(r"uddg=([^&]+)", href)
        if redirect:
            url = unquote(redirect.group(1))
        title = unescape(re.sub(r"<[^>]+>", "", raw_title)).strip()
        snippet = ""
        if index < len(snippets):
            snippet = unescape(re.sub(r"<[^>]+>", "", snippets[index])).strip()
        results.append(f"{index + 1}. {title}\n   {url}\n   {snippet}")

    if not results:
        return f"No search results found for: {query}"
    return "Search results:\n\n" + "\n\n".join(results)


def fetch_webpage(db: Session, chat_id: str, url: str):
    """Read a page headlessly, save its text to chat memory, return the text."""
    if not re.match(r"^https?://", url, re.I):
        url = f"https://{url}"

    page = fetch_page(url, headless=True, wait_seconds=2)
    chunk_count = save_document_text(db, chat_id, page.title, f"URL: {page.url}\n\n{page.text}")

    parts = [f"Page: {page.title} ({page.url})", f"Saved {chunk_count} chunks to chat memory.", "", clip(page.text)]
    if len(page.text.split()) < 15:
        parts.append(
            "\nNote: the page returned almost no text. It may require login or heavy JavaScript — "
            "consider open_browser so the user can log in, then read_browser_page."
        )
    if page.pdf_links:
        parts.append("\nPDF links on the page (use read_pdf to open one):\n" + "\n".join(page.pdf_links[:10]))
    if page.links:
        listed = "\n".join(f"- [{link['type']}] {link['label']}: {link['url']}" for link in page.links[:20])
        parts.append("\nLinks on the page:\n" + listed)
    return "\n".join(parts)


def read_pdf_url(db: Session, chat_id: str, url: str):
    response = requests.get(url, timeout=30)
    response.raise_for_status()
    text = extract_file_text(url, response.content)
    if not text.strip():
        return "The PDF downloaded but no text could be extracted (it may be scanned images)."
    chunk_count = save_document_text(db, chat_id, url, text)
    return f"PDF read and saved to chat memory ({chunk_count} chunks).\n\n{clip(text)}"


# ---------------------------------------------------------------------------
# Visible browser sessions (Selenium Chrome the user can see and log into)
# ---------------------------------------------------------------------------

def _browser_session(chat_id: str):
    return ACTIVE_BROWSER_SESSIONS.get(str(chat_id))


def open_browser(chat_id: str, url: str):
    if not re.match(r"^https?://", url, re.I):
        if re.match(r"^(?:www\.)?[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}(?:/\S*)?$", url.strip()):
            url = f"https://{url.strip()}"
        else:
            url = f"https://www.google.com/search?q={quote_plus(url)}"

    close_browser(chat_id)
    driver = create_driver(headless=False)
    try:
        driver.get(url)
    except Exception as exc:
        driver.quit()
        raise RuntimeError(f"Could not open {url}: {exc}") from exc

    ACTIVE_BROWSER_SESSIONS[str(chat_id)] = {"driver": driver, "url": url}
    return (
        f"Opened a visible Chrome window at {url}. The user can interact with it "
        "(log in, navigate). Use read_browser_page to read whatever is currently shown."
    )


def close_browser(chat_id: str):
    session = ACTIVE_BROWSER_SESSIONS.pop(str(chat_id), None)
    if not session:
        return "No browser window is open for this chat."
    try:
        session["driver"].quit()
    except Exception:
        pass
    return "Closed the browser window."


def read_browser_page(db: Session, chat_id: str):
    session = _browser_session(chat_id)
    if not session:
        return "No browser window is open. Use open_browser first (or fetch_webpage for pages that need no login)."

    driver = session["driver"]
    try:
        title = driver.title or "Open page"
        url = driver.current_url
        body_text = driver.find_element(By.TAG_NAME, "body").text.strip()
    except Exception as exc:
        raise RuntimeError(f"Could not read the browser page — the window may have been closed: {exc}") from exc

    if len(body_text.split()) < 5:
        return "The open page has almost no readable text. The user may still be logging in or the page is still loading."

    links = []
    pdf_links = []
    for link in driver.find_elements(By.CSS_SELECTOR, "a[href]")[:120]:
        href = link.get_attribute("href")
        if not href:
            continue
        label = (link.text.strip() or link.get_attribute("aria-label") or href)[:120]
        if ".pdf" in href.lower():
            pdf_links.append(href)
        links.append(f"- {label}: {href}")

    chunk_count = save_document_text(db, chat_id, title, f"URL: {url}\n\n{body_text}")

    parts = [f"Current page: {title} ({url})", f"Saved {chunk_count} chunks to chat memory.", "", clip(body_text)]
    if pdf_links:
        parts.append("\nPDF links (use read_pdf):\n" + "\n".join(sorted(set(pdf_links))[:10]))
    if links:
        parts.append("\nLinks:\n" + "\n".join(links[:20]))
    return "\n".join(parts)


def browser_navigate(chat_id: str, url: str):
    session = _browser_session(chat_id)
    if not session:
        return open_browser(chat_id, url)
    if not re.match(r"^https?://", url, re.I):
        url = f"https://{url}"
    session["driver"].get(url)
    session["url"] = url
    return f"Navigated the open browser to {url}. Use read_browser_page to read it."


def read_quiz_questions(chat_id: str):
    session = _browser_session(chat_id)
    if not session:
        return "No browser window is open. Use open_browser with the quiz URL first, let the user log in, then call this again."

    questions = extract_mcq_questions(session["driver"])
    if not questions:
        return "No visible radio-button or checkbox questions found on the current page."

    session["questions"] = questions
    blocks = []
    for q_index, question in enumerate(questions, start=1):
        options = "\n".join(
            f"  {chr(64 + option_index)}. {option.label}"
            for option_index, option in enumerate(question.options, start=1)
        )
        blocks.append(f"Q{q_index}. {question.question}\n{options}")
    return (
        "Visible questions on the page (answer choices are labeled A, B, C...):\n\n"
        + "\n\n".join(blocks)
        + "\n\nTo click choices use select_quiz_answers with e.g. "
        '[{"question": 1, "choice": "B"}]. Selections are clicked but the quiz is never submitted.'
    )


def select_quiz_answers(chat_id: str, answers):
    session = _browser_session(chat_id)
    if not session or not session.get("questions"):
        return "Call read_quiz_questions first so I know which options exist on the page."

    driver = session["driver"]
    questions = session["questions"]
    clicked, skipped = [], []

    for item in answers:
        try:
            q_num = int(item.get("question"))
            choice = str(item.get("choice", "")).strip().upper()[:1]
            option_index = ord(choice) - 65
            option = questions[q_num - 1].options[option_index]
        except (ValueError, TypeError, IndexError, AttributeError):
            skipped.append(str(item))
            continue

        try:
            element = next(
                input_el
                for input_el in driver.find_elements(By.CSS_SELECTOR, "input[type='radio'], input[type='checkbox']")
                if input_el.id == option.element_id
            )
            driver.execute_script("arguments[0].scrollIntoView({block: 'center'});", element)
            element.click()
            clicked.append(f"Q{q_num}={choice}")
        except Exception:
            skipped.append(f"Q{q_num}")

    return (
        f"Clicked: {', '.join(clicked) or 'none'}. Skipped: {', '.join(skipped) or 'none'}. "
        "The quiz was NOT submitted — the user must review and submit themselves."
    )


# ---------------------------------------------------------------------------
# Desktop control: open apps, keyboard, typing
# ---------------------------------------------------------------------------

def open_local_target(target: str):
    if os.name == "nt":
        if re.match(r"^https?://", target, re.I) or Path(target).exists():
            os.startfile(target)  # type: ignore[attr-defined]
        else:
            subprocess.run(
                ["powershell", "-NoProfile", "-Command", "Start-Process", "-FilePath", target],
                check=True,
                capture_output=True,
                text=True,
            )
        return
    subprocess.Popen(["xdg-open", target])


def _find_windows_app_path(exe_name: str):
    if os.name != "nt" or not exe_name:
        return None
    try:
        import winreg
    except ImportError:
        return None

    key_paths = [
        (winreg.HKEY_LOCAL_MACHINE, fr"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\{exe_name}.exe"),
        (winreg.HKEY_CURRENT_USER, fr"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\{exe_name}.exe"),
        (winreg.HKEY_LOCAL_MACHINE, fr"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\App Paths\{exe_name}.exe"),
    ]
    for hive, key_path in key_paths:
        try:
            with winreg.OpenKey(hive, key_path) as key:
                value, _ = winreg.QueryValueEx(key, None)
                if value and Path(value).exists():
                    return value
        except OSError:
            continue
    return None


def _find_common_windows_executable(exe_name: str):
    if os.name != "nt" or not exe_name:
        return None
    candidates = [
        Path(os.environ.get("ProgramFiles", "")) / "Microsoft Office" / "root" / "Office16" / exe_name,
        Path(os.environ.get("ProgramFiles(x86)", "")) / "Microsoft Office" / "root" / "Office16" / exe_name,
        Path(os.environ.get("LOCALAPPDATA", "")) / "Programs" / "Microsoft VS Code" / "Code.exe",
        Path(os.environ.get("ProgramFiles", "")) / "Google" / "Chrome" / "Application" / "chrome.exe",
        Path(os.environ.get("ProgramFiles(x86)", "")) / "Microsoft" / "Edge" / "Application" / "msedge.exe",
    ]
    for candidate in candidates:
        if candidate.name.lower() == exe_name.lower() and candidate.exists():
            return str(candidate)
    return None


def open_app(chat_id: str, name: str):
    lowered = name.lower().strip()
    registry = app_registry()

    matched_label, info = None, None
    for label, entry in registry.items():
        if label in lowered:
            matched_label, info = label, entry
            break

    if info:
        command = info.get("command")
        attempts = [
            _find_windows_app_path(command) if command else None,
            _find_common_windows_executable(info.get("exe")) if info.get("exe") else None,
            command,
            info.get("uri"),
        ]
        last_error = None
        for target in [item for item in attempts if item]:
            try:
                open_local_target(target)
                ACTIVE_APP_SESSIONS[str(chat_id)] = matched_label
                return f"Opened {matched_label.title()}. It should now be the focused window."
            except Exception as exc:
                last_error = exc
        raise RuntimeError(f"Could not open {matched_label}: {last_error}")

    # Unknown app: let the OS try to resolve the name directly.
    try:
        open_local_target(lowered)
        ACTIVE_APP_SESSIONS[str(chat_id)] = lowered
        return f"Asked Windows to open '{name}'."
    except Exception as exc:
        raise RuntimeError(
            f"Windows could not find an app named '{name}'. Known apps: {', '.join(sorted(registry))}."
        ) from exc


def send_keys(keys: str):
    """Send SendKeys-format keystrokes to the focused window (e.g. ^s, {ENTER}, %{F4})."""
    if os.name != "nt":
        raise RuntimeError("Keyboard control is only implemented for Windows.")

    subprocess.run(
        [
            "powershell",
            "-NoProfile",
            "-Command",
            (
                "Add-Type -AssemblyName System.Windows.Forms; "
                "Start-Sleep -Milliseconds 350; "
                f"[System.Windows.Forms.SendKeys]::SendWait('{keys}')"
            ),
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    return f"Sent keystrokes {keys} to the focused window."


def type_text(text: str):
    """Paste text into the focused window via clipboard."""
    if os.name != "nt":
        raise RuntimeError("Typing into apps is only implemented for Windows.")

    subprocess.run(
        [
            "powershell",
            "-NoProfile",
            "-Command",
            (
                "Add-Type -AssemblyName System.Windows.Forms; "
                "$text = [Console]::In.ReadToEnd(); "
                "Set-Clipboard -Value $text; "
                "Start-Sleep -Milliseconds 250; "
                "[System.Windows.Forms.SendKeys]::SendWait('^v')"
            ),
        ],
        input=text,
        text=True,
        check=True,
        capture_output=True,
    )
    return "Pasted the text into the focused window. The user should review it before saving or submitting."


# ---------------------------------------------------------------------------
# File generation
# ---------------------------------------------------------------------------

def safe_filename(name: str, fallback="assistant_output"):
    cleaned = re.sub(r"[^a-zA-Z0-9_. -]+", "", name).strip().strip(".")
    cleaned = re.sub(r"\s+", "_", cleaned)
    return cleaned[:80] or fallback


def generated_file_url(path: Path):
    return f"/generated/{path.name}"


def create_file(file_type: str, title: str, content: str):
    """Create a docx/xlsx/pptx/txt file with the given content and return its path + URL."""
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    file_type = (file_type or "docx").lower().lstrip(".")
    filename = f"{safe_filename(title)}_{uuid4().hex[:8]}.{file_type}"
    path = GENERATED_DIR / filename

    if file_type == "docx":
        from docx import Document

        document = Document()
        document.add_heading(title, level=1)
        for block in content.split("\n\n"):
            if block.strip():
                document.add_paragraph(block.strip())
        document.save(path)
    elif file_type == "xlsx":
        from openpyxl import Workbook

        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Sheet1"
        for line in content.splitlines():
            if line.strip():
                cells = [cell.strip() for cell in re.split(r"\t|\||,", line)]
                sheet.append(cells)
        workbook.save(path)
    elif file_type == "pptx":
        from pptx import Presentation

        presentation = Presentation()
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = title
        blocks = [block.strip() for block in content.split("\n\n") if block.strip()]
        for block in blocks[:12]:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            lines = block.splitlines()
            slide.shapes.title.text = lines[0][:80]
            slide.placeholders[1].text = "\n".join(lines[1:])[:900] or block[:900]
        presentation.save(path)
    elif file_type == "txt":
        path.write_text(content, encoding="utf-8")
    else:
        raise ValueError(f"Unsupported file type '{file_type}'. Use docx, xlsx, pptx, or txt.")

    return f"Created {path.name}. Local path: {path}. Download link: {generated_file_url(path)}"
