from uuid import uuid4
import base64
import io
import json
import os
import re
import secrets
import subprocess
from datetime import datetime, timedelta
from pathlib import Path
from urllib.parse import urlencode

import pdfplumber
import pytesseract
import requests
from dotenv import load_dotenv
from fastapi import FastAPI, UploadFile, File, Form, Depends, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, HTMLResponse, RedirectResponse, StreamingResponse
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel
from PIL import Image
from sqlalchemy.orm import Session
from selenium.webdriver.common.by import By

from auth import (
    create_access_token,
    create_refresh_token,
    decode_token,
    hash_password,
    verify_password,
)
from database import Base, SessionLocal, engine
from models import Chat, DocumentChunk, ImageAttachment, Message, User
from scraper import create_driver, extract_mcq_questions, scrape_assignment_page
from vector_store import retrieve_relevant_chunks, split_text


load_dotenv()

app = FastAPI()
security = HTTPBearer()
optional_security = HTTPBearer(auto_error=False)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://127.0.0.1:5173", "http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

Base.metadata.create_all(bind=engine)

OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://127.0.0.1:11434")
OLLAMA_URL = os.getenv("OLLAMA_URL", f"{OLLAMA_BASE_URL}/api/generate")
OLLAMA_TIMEOUT_SECONDS = int(os.getenv("OLLAMA_TIMEOUT_SECONDS", "300"))
DEFAULT_OLLAMA_MODEL = os.getenv("OLLAMA_MODEL", "phi3:mini")
VISION_MODEL_KEYWORDS = (
    "llava",
    "vision",
    "bakllava",
    "moondream",
    "minicpm-v",
    "qwen-vl",
    "qwen2-vl",
    "qwen2.5vl",
    "qwen2.5-vl",
)
VISION_OLLAMA_MODELS = {"llava:latest", "llama3.2-vision:latest"}
ALLOWED_OLLAMA_MODELS = {
    "phi3:mini",
    "llama3:latest",
    "neural-chat:latest",
    "mistral:latest",
    *VISION_OLLAMA_MODELS
}
MODEL_ALIASES = {
    "llama3": "llama3:latest",
    "llama3.1": "llama3.1:latest",
    "llava": "llava:latest",
    "llama3.2-vision": "llama3.2-vision:latest",
}
UPLOAD_DIR = Path(__file__).resolve().parent / "uploads"
GENERATED_DIR = Path(__file__).resolve().parent / "generated"
ACTIVE_MCQ_SESSIONS = {}
ACTIVE_BROWSER_SESSIONS = {}
ACTIVE_APP_SESSIONS = {}
PASSWORD_RESET_TOKENS = {}
PASSWORD_RESET_CODES = {}
OAUTH_STATES = {}
PASSWORD_RESET_EXPIRE_MINUTES = 20
RESET_CODE_EXPIRE_MINUTES = 10
OAUTH_STATE_EXPIRE_MINUTES = 10
FRONTEND_URL = os.getenv("FRONTEND_URL", "http://127.0.0.1:5173")
BACKEND_PUBLIC_URL = os.getenv("BACKEND_PUBLIC_URL", "http://127.0.0.1:8000")
OAUTH_PROVIDERS = {
    "google": {
        "client_id_env": "GOOGLE_CLIENT_ID",
        "client_secret_env": "GOOGLE_CLIENT_SECRET",
        "authorize_url": "https://accounts.google.com/o/oauth2/v2/auth",
        "token_url": "https://oauth2.googleapis.com/token",
        "userinfo_url": "https://openidconnect.googleapis.com/v1/userinfo",
        "scope": "openid email profile",
    },
}
DEFAULT_WEB_TARGETS = {
    "canvas": "https://canvas.instructure.com/",
    "instructure": "https://canvas.instructure.com/",
    "google classroom": "https://classroom.google.com/",
    "classroom": "https://classroom.google.com/",
    "blackboard": "https://www.blackboard.com/",
    "moodle": "https://moodle.org/login/",
}
DEFAULT_APP_REGISTRY = {
    "microsoft word": {"command": "winword", "exe": "WINWORD.EXE", "category": "document"},
    "word": {"command": "winword", "exe": "WINWORD.EXE", "category": "document"},
    "microsoft excel": {"command": "excel", "exe": "EXCEL.EXE", "category": "spreadsheet"},
    "excel": {"command": "excel", "exe": "EXCEL.EXE", "category": "spreadsheet"},
    "powerpoint": {"command": "powerpnt", "exe": "POWERPNT.EXE", "category": "presentation"},
    "power point": {"command": "powerpnt", "exe": "POWERPNT.EXE", "category": "presentation"},
    "chrome": {"command": "chrome", "category": "browser"},
    "browser": {"command": "chrome", "category": "browser"},
    "edge": {"command": "msedge", "category": "browser"},
    "notepad": {"command": "notepad", "category": "document"},
    "calculator": {"command": "calc", "category": "calculator"},
    "calc": {"command": "calc", "category": "calculator"},
    "camera": {"uri": "microsoft.windows.camera:", "category": "camera"},
    "photos": {"uri": "ms-photos:", "category": "media"},
    "settings": {"uri": "ms-settings:", "category": "system"},
    "mail": {"uri": "outlookmail:", "category": "communication"},
    "calendar": {"uri": "outlookcal:", "category": "calendar"},
    "store": {"uri": "ms-windows-store:", "category": "store"},
    "vscode": {"command": "code", "category": "editor"},
    "vs code": {"command": "code", "category": "editor"},
}
GENERIC_KEY_ACTIONS = {
    "save": "^s",
    "print": "^p",
    "select all": "^a",
    "copy": "^c",
    "cut": "^x",
    "paste clipboard": "^v",
    "undo": "^z",
    "redo": "^y",
    "bold": "^b",
    "italic": "^i",
    "underline": "^u",
    "tab": "{TAB}",
    "escape": "{ESC}",
    "cancel": "{ESC}",
    "press enter": "{ENTER}",
    "next line": "{ENTER}",
    "new line": "{ENTER}",
    "new file": "^n",
    "new document": "^n",
    "new workbook": "^n",
    "new presentation": "^n",
    "blank document": "^n",
    "blank workbook": "^n",
    "blank presentation": "^n",
    "close app": "%{F4}",
    "close it": "%{F4}",
    "close window": "%{F4}",
}


class AuthRequest(BaseModel):
    email: str
    password: str


class ForgotPasswordRequest(BaseModel):
    email: str


class ResetPasswordRequest(BaseModel):
    token: str
    password: str


class ResetCodeRequest(BaseModel):
    email: str
    code: str


class ChangePasswordRequest(BaseModel):
    email: str
    reset_token: str
    password: str


class SocialLoginRequest(BaseModel):
    email: str
    provider: str = "google"


class StreamRequest(BaseModel):
    chat_id: str
    message: str
    model: str | None = None


class RenameRequest(BaseModel):
    title: str


class ScrapeRequest(BaseModel):
    url: str
    headless: bool = True
    wait_seconds: int = 2
    include_pdfs: bool = True
    manual_login: bool = False
    login_wait_seconds: int = 90


class ActionRequest(BaseModel):
    chat_id: str
    message: str


def generated_file_url(path: Path):
    return f"/generated/{path.name}"


def resolve_generated_file(filename: str):
    safe_name = Path(filename).name
    path = (GENERATED_DIR / safe_name).resolve()
    root = GENERATED_DIR.resolve()
    if root not in path.parents or not path.exists() or not path.is_file():
        raise HTTPException(status_code=404, detail="Generated file not found")
    return path


def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


def mask_email(email: str):
    name, _, domain = email.partition("@")
    if len(name) <= 2:
        masked_name = name[:1] + "*"
    else:
        masked_name = f"{name[:2]}{'*' * min(5, len(name) - 2)}"
    return f"{masked_name}@{domain}"


def email_configured():
    return bool(os.getenv("RESEND_API_KEY") and os.getenv("EMAIL_FROM"))


def reset_code_email_content(code: str):
    text = (
        f"Use this code to reset your Study Assistant password: {code}\n\n"
        f"This code expires in {RESET_CODE_EXPIRE_MINUTES} minutes. If you did not request it, you can ignore this email."
    )
    html = f"""
        <div style="font-family:Arial,sans-serif;max-width:560px;color:#111827">
          <h2>Reset your password</h2>
          <p>Use this code to reset your Study Assistant password:</p>
          <div style="font-size:28px;font-weight:700;letter-spacing:6px;padding:14px 18px;background:#f3f4f6;border-radius:8px;display:inline-block">{code}</div>
          <p style="color:#667085">This code expires in {RESET_CODE_EXPIRE_MINUTES} minutes. If you did not request it, you can ignore this email.</p>
        </div>
        """
    return text, html


def send_reset_code_resend(email: str, code: str):
    if not email_configured():
        raise RuntimeError("Resend email is not configured. Set RESEND_API_KEY and EMAIL_FROM.")

    text, html = reset_code_email_content(code)
    response = requests.post(
        "https://api.resend.com/emails",
        headers={
            "Authorization": f"Bearer {os.getenv('RESEND_API_KEY')}",
            "Content-Type": "application/json",
        },
        json={
            "from": os.getenv("EMAIL_FROM"),
            "to": [email],
            "subject": "Your Study Assistant password reset code",
            "text": text,
            "html": html,
        },
        timeout=20,
    )
    response.raise_for_status()


def send_reset_code_email(email: str, code: str):
    if not email_configured():
        raise RuntimeError(
            "Email is not configured. Set RESEND_API_KEY and EMAIL_FROM."
        )
    send_reset_code_resend(email, code)


def create_reset_code(email: str):
    code = "".join(secrets.choice("0123456789") for _ in range(6))
    PASSWORD_RESET_CODES[email] = {
        "code": code,
        "expires_at": datetime.utcnow() + timedelta(minutes=RESET_CODE_EXPIRE_MINUTES),
        "attempts": 0,
        "verified": False,
    }
    return code


def oauth_redirect_uri(provider: str):
    return f"{BACKEND_PUBLIC_URL}/oauth/{provider}/callback"


def get_or_create_oauth_user(db: Session, email: str):
    db_user = db.query(User).filter(User.email == email).first()
    if not db_user:
        db_user = User(id=str(uuid4()), email=email, password=hash_password(secrets.token_urlsafe(24)))
        db.add(db_user)
        db.commit()
        db.refresh(db_user)
    return db_user


def oauth_success_page(access_token: str, refresh_token: str):
    return HTMLResponse(f"""
    <!doctype html>
    <html>
      <body>
        <script>
          localStorage.setItem("token", {json.dumps(access_token)});
          localStorage.setItem("refresh_token", {json.dumps(refresh_token)});
          window.location.replace({json.dumps(FRONTEND_URL)});
        </script>
        Signing you in...
      </body>
    </html>
    """)


def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)):
    data = decode_token(credentials.credentials)
    if not data or data.get("type") != "access":
        raise HTTPException(status_code=401, detail="Unauthorized")
    return data["user"]


def get_current_user_optional(
    token: str | None = None,
    credentials: HTTPAuthorizationCredentials | None = Depends(optional_security),
):
    raw_token = token or (credentials.credentials if credentials else None)
    data = decode_token(raw_token) if raw_token else None
    if not data or data.get("type") != "access":
        raise HTTPException(status_code=401, detail="Unauthorized")
    return data["user"]


def get_owned_chat(chat_id: str, user: str, db: Session):
    chat = db.query(Chat).filter(Chat.id == chat_id, Chat.user_id == user).first()
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    return chat


def extract_file_text(filename: str, file_bytes: bytes):
    lower_name = filename.lower()

    if lower_name.endswith(".pdf"):
        content = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    content.append(text)
        return "\n".join(content)

    if lower_name.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff")):
        image = Image.open(io.BytesIO(file_bytes))
        return pytesseract.image_to_string(image)

    if lower_name.endswith(".docx"):
        from docx import Document

        document = Document(io.BytesIO(file_bytes))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)

    if lower_name.endswith((".xlsx", ".xlsm")):
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
        rows = []
        for sheet in workbook.worksheets:
            rows.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value) for value in row if value is not None]
                if values:
                    rows.append(" | ".join(values))
        return "\n".join(rows)

    if lower_name.endswith((".doc", ".xls")):
        raise HTTPException(
            status_code=415,
            detail="Old .doc/.xls files are not supported yet. Save as .docx or .xlsx and upload again.",
        )

    return file_bytes.decode("utf-8", errors="replace")


def is_image_filename(filename: str):
    return filename.lower().endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"))


def save_image_attachment(db: Session, chat_id: str, filename: str, file_bytes: bytes):
    suffix = Path(filename).suffix.lower() or ".jpg"
    image_dir = UPLOAD_DIR / chat_id
    image_dir.mkdir(parents=True, exist_ok=True)
    image_path = image_dir / f"{uuid4()}{suffix}"
    image_path.write_bytes(file_bytes)

    db.add(ImageAttachment(
        id=str(uuid4()),
        chat_id=chat_id,
        filename=filename,
        path=str(image_path),
    ))
    return image_path


def get_latest_image_payloads(db: Session, chat_id: str, limit=1):
    rows = db.query(ImageAttachment).filter(ImageAttachment.chat_id == chat_id).all()
    payloads = []
    for row in rows[-limit:]:
        image_path = Path(row.path)
        if image_path.exists():
            payloads.append(base64.b64encode(image_path.read_bytes()).decode("utf-8"))
    return payloads


def extract_pdf_url_text(url: str):
    response = requests.get(url, timeout=30)
    response.raise_for_status()
    return extract_file_text(url, response.content)


def build_context(db: Session, chat_id: str, query: str):
    chunks = [
        f"Source file: {row.filename}\nExtracted text:\n{row.content}"
        for row in db.query(DocumentChunk).filter(DocumentChunk.chat_id == chat_id).all()
    ]
    return "\n\n".join(retrieve_relevant_chunks(chunks, query))


def save_document_text(db: Session, chat_id: str, filename: str, content: str):
    chunks = split_text(content.strip())
    if not chunks:
        return 0

    db.add_all([
        DocumentChunk(
            id=str(uuid4()),
            chat_id=chat_id,
            filename=filename,
            content=chunk,
        )
        for chunk in chunks
    ])
    return len(chunks)


def safe_filename(name: str, fallback="jarvis_output"):
    cleaned = re.sub(r"[^a-zA-Z0-9_. -]+", "", name).strip().strip(".")
    cleaned = re.sub(r"\s+", "_", cleaned)
    return cleaned[:80] or fallback


def open_local_target(target: str):
    if os.name == "nt":
        if re.match(r"^https?://", target, re.I) or Path(target).exists():
            os.startfile(target)  # type: ignore[attr-defined]
        else:
            subprocess.run(
                ["powershell", "-NoProfile", "-Command", "Start-Process", "-FilePath", target],
                check=True,
                capture_output=True,
                text=True,
            )
        return
    subprocess.Popen(["xdg-open", target])


def open_website(url: str):
    if not re.match(r"^https?://", url, re.I):
        url = f"https://{url}"
    open_local_target(url)
    return url


def extract_search_query(message: str):
    text = message.strip()
    patterns = [
        r"(?is)^\s*search\s+(?:for\s+)?(.+?)(?:\s+on\s+(?:it|chrome|browser|google))?\s*$",
        r"(?is)^\s*google\s+(.+?)\s*$",
        r"(?is)^\s*look\s+up\s+(.+?)\s*$",
    ]
    for pattern in patterns:
        match = re.match(pattern, text)
        if match:
            query = match.group(1).strip(" .")
            if query:
                return query
    return None


def open_browser_search(query: str):
    from urllib.parse import quote_plus

    url = f"https://www.google.com/search?q={quote_plus(query)}"
    open_local_target(url)
    return url


def load_json_mapping_env(name: str, fallback: dict):
    raw = os.getenv(name)
    if not raw:
        return dict(fallback)
    try:
        loaded = json.loads(raw)
        if isinstance(loaded, dict):
            merged = dict(fallback)
            merged.update(loaded)
            return merged
    except json.JSONDecodeError:
        pass
    return dict(fallback)


def web_targets():
    return load_json_mapping_env("JARVIS_WEB_TARGETS", DEFAULT_WEB_TARGETS)


def app_registry():
    return load_json_mapping_env("JARVIS_APP_REGISTRY", DEFAULT_APP_REGISTRY)


def extract_known_web_app(message: str):
    lowered = message.lower()
    for label, url in web_targets().items():
        if re.search(rf"\b{re.escape(label)}\b", lowered):
            return label, url
    return None, None


def extract_browser_open_target(message: str):
    text = message.strip()
    if not re.match(r"(?is)^\s*open\b", text):
        return None

    target = re.sub(r"(?is)^\s*open\s+", "", text).strip()
    target = re.sub(r"(?is)\b(?:in|on|with)\s+(?:chrome|browser|web\s*browser)\b", "", target).strip()
    target = re.sub(r"(?is)\b(?:chrome|browser|web\s*browser)\b", "", target).strip()
    target = re.sub(r"(?is)\b(?:please|for me|the|a)\b", "", target).strip(" .")
    return target or None


def browser_url_for_target(target: str):
    from urllib.parse import quote_plus

    cleaned = target.strip()
    if re.match(r"^https?://", cleaned, re.I):
        return cleaned
    if re.match(r"^(?:www\.)?[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}(?:/\S*)?$", cleaned):
        return f"https://{cleaned}"
    return f"https://www.google.com/search?q={quote_plus(cleaned)}"


def looks_like_web_target(target: str):
    cleaned = target.strip().lower()
    return bool(
        re.match(r"^https?://", cleaned)
        or re.match(r"^(?:www\.)?[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}(?:/\S*)?$", cleaned)
        or any(word in cleaned for word in ["website", "webpage", "web page", "site", "url", "link"])
    )


def allowed_app_command(app_name: str):
    lowered = app_name.lower().strip()
    for label, command in app_registry().items():
        if label in lowered:
            return label, command
    return None, None


def extract_open_app_name(message: str):
    match = re.match(r"(?i)^\s*open\s+(.+?)\s*$", message)
    if not match:
        return None

    app_name = match.group(1).strip()
    if re.match(r"^https?://", app_name, re.I) or "." in app_name:
        return None
    app_name = re.sub(r"\b(app|application|please|for me)\b", "", app_name, flags=re.I).strip()
    if not app_name:
        return None
    if not re.match(r"^[a-zA-Z0-9 ._-]{2,60}$", app_name):
        return None
    return app_name


def find_windows_app_path(exe_name: str):
    if os.name != "nt":
        return None

    try:
        import winreg
    except ImportError:
        return None

    key_paths = [
        (winreg.HKEY_LOCAL_MACHINE, fr"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\{exe_name}.exe"),
        (winreg.HKEY_CURRENT_USER, fr"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\{exe_name}.exe"),
        (winreg.HKEY_LOCAL_MACHINE, fr"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\App Paths\{exe_name}.exe"),
    ]
    for hive, key_path in key_paths:
        try:
            with winreg.OpenKey(hive, key_path) as key:
                value, _ = winreg.QueryValueEx(key, None)
                if value and Path(value).exists():
                    return value
        except OSError:
            continue
    return None


def find_common_windows_executable(exe_name: str):
    if os.name != "nt" or not exe_name:
        return None

    candidates = [
        Path(os.environ.get("ProgramFiles", "")) / "Microsoft Office" / "root" / "Office16" / exe_name,
        Path(os.environ.get("ProgramFiles(x86)", "")) / "Microsoft Office" / "root" / "Office16" / exe_name,
        Path(os.environ.get("LOCALAPPDATA", "")) / "Programs" / "Microsoft VS Code" / "Code.exe",
        Path(os.environ.get("ProgramFiles", "")) / "Google" / "Chrome" / "Application" / "chrome.exe",
        Path(os.environ.get("ProgramFiles(x86)", "")) / "Microsoft" / "Edge" / "Application" / "msedge.exe",
    ]
    for candidate in candidates:
        if candidate.name.lower() == exe_name.lower() and candidate.exists():
            return str(candidate)

    return None


def launch_allowed_app(app_info: dict):
    command = app_info.get("command")
    exe = app_info.get("exe")
    uri = app_info.get("uri")

    if os.name == "nt":
        app_path = find_windows_app_path(command) if command else None
        common_path = find_common_windows_executable(exe) if exe else None
        attempts = [app_path, common_path, command, uri]
        last_error = None
        for target in [item for item in attempts if item]:
            try:
                open_local_target(target)
                return
            except Exception as exc:
                last_error = exc
        raise RuntimeError(f"Windows could not open that app. It may not be installed. Details: {last_error}")

    if command:
        subprocess.Popen([command])
        return
    raise RuntimeError("I do not know how to open that app on this operating system.")


def launch_named_app(app_name: str):
    if os.name == "nt":
        try:
            open_local_target(app_name)
            return
        except Exception as exc:
            raise RuntimeError(
                f"Windows could not find an app named `{app_name}`. Try the exact app name, like `open chrome`, `open word`, or `open notepad`."
            ) from exc

    subprocess.Popen([app_name])


def app_session_key(chat_id: str):
    return str(chat_id)


def remember_active_app(chat_id: str, app_label: str, app_info: dict | None = None):
    ACTIVE_APP_SESSIONS[app_session_key(chat_id)] = {
        "label": app_label,
        "info": app_info or {},
    }


def latest_active_app(chat_id: str):
    return ACTIVE_APP_SESSIONS.get(app_session_key(chat_id))


def active_app_label(chat_id: str):
    active = latest_active_app(chat_id)
    return (active.get("label") if active else "").lower()


def active_app_category(chat_id: str):
    active = latest_active_app(chat_id)
    info = active.get("info", {}) if active else {}
    return info.get("category") or (active.get("label", "") if active else "")


def active_app_display_name(chat_id: str):
    active = latest_active_app(chat_id)
    return active.get("label", "the active app").title() if active else "the active app"


def send_keys_to_active_window(keys: str):
    if os.name != "nt":
        raise RuntimeError("Controlling the active app is currently implemented for Windows.")

    subprocess.run(
        [
            "powershell",
            "-NoProfile",
            "-Command",
            (
                "Add-Type -AssemblyName System.Windows.Forms; "
                "Start-Sleep -Milliseconds 350; "
                f"[System.Windows.Forms.SendKeys]::SendWait('{keys}')"
            ),
        ],
        check=True,
        capture_output=True,
        text=True,
    )


def wants_calculator_action(lowered: str, chat_id: str):
    group = active_app_category(chat_id)
    has_expression = re.search(r"[-+*/xX÷×=().\d ]{3,}", lowered)
    return group == "calculator" and (
        has_expression
        or any(term in lowered for term in ["calculate", "equals", "press enter", "clear", "backspace"])
    )


def extract_calculator_expression(message: str):
    text = message.lower()
    text = re.sub(r"\b(?:calculate|calculator|what is|what's|equals|equal to|enter|type|please)\b", " ", text)
    text = text.replace("plus", "+").replace("minus", "-").replace("times", "*")
    text = text.replace("multiplied by", "*").replace("divided by", "/")
    text = text.replace("x", "*").replace("×", "*").replace("÷", "/")
    allowed = re.findall(r"[\d+\-*/().= ]+", text)
    expression = " ".join(allowed).strip()
    expression = re.sub(r"\s+", "", expression).rstrip("=")
    return expression if re.search(r"\d", expression) else None


def run_calculator_action(chat_id: str, message: str):
    lowered = message.lower()
    if "clear" in lowered:
        send_keys_to_active_window("{ESC}")
        return True, "Cleared Calculator."
    if "backspace" in lowered or "delete last" in lowered:
        send_keys_to_active_window("{BACKSPACE}")
        return True, "Sent Backspace to Calculator."

    expression = extract_calculator_expression(message)
    if not expression:
        return True, "Tell me the calculation, like `calculate 25 * 4`."

    paste_text_into_active_window(expression)
    send_keys_to_active_window("{ENTER}")
    return True, f"Entered `{expression}` in Calculator."


def wants_generic_key_action(lowered: str, chat_id: str):
    if not latest_active_app(chat_id):
        return False
    return any(term in lowered for term in GENERIC_KEY_ACTIONS)


def run_generic_key_action(chat_id: str, message: str):
    lowered = message.lower()
    target_label = active_app_display_name(chat_id)
    for phrase, keys in sorted(GENERIC_KEY_ACTIONS.items(), key=lambda item: len(item[0]), reverse=True):
        if phrase in lowered:
            send_keys_to_active_window(keys)
            return True, f"Sent `{phrase}` to {target_label}."

    return False, ""


def wants_camera_capture(lowered: str, chat_id: str):
    active = latest_active_app(chat_id)
    camera_is_active = active_app_category(chat_id) == "camera" if active else False
    capture_terms = ["take picture", "take photo", "click picture", "capture picture", "capture photo", "snap picture", "snap photo"]
    return any(term in lowered for term in capture_terms) and (camera_is_active or "camera" in lowered)


def capture_camera_photo(chat_id: str):
    active = latest_active_app(chat_id)
    if not active or active_app_category(chat_id) != "camera":
        app_label, app_command = allowed_app_command("camera")
        launch_allowed_app(app_command)
        remember_active_app(chat_id, app_label, app_command)

    send_keys_to_active_window("{ENTER}")
    return True, "I sent the camera shutter command. If Camera asked for permission or focus moved, click the Camera window once and say `click picture` again."


def make_generation_text(message: str, context: str):
    context_text = make_preview(context, 2200) if context else ""
    if context_text:
        return (
            f"Request\n{message.strip()}\n\n"
            f"Relevant chat or assignment context\n{context_text}\n\n"
            "Draft\nUse this file as a starting point. Ask in chat for a more polished version, citations, code, or formatting changes."
        )
    return (
        f"Request\n{message.strip()}\n\n"
        "Draft\nUse this file as a starting point. Ask in chat for a more polished version, citations, code, or formatting changes."
    )


def create_docx_file(path: Path, title: str, body: str):
    from docx import Document

    document = Document()
    document.add_heading(title, level=1)
    for block in body.split("\n\n"):
        if block.strip():
            document.add_paragraph(block.strip())
    document.save(path)


def create_xlsx_file(path: Path, title: str, body: str):
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Assistant Draft"
    sheet.append(["Section", "Content"])
    sheet.append(["Title", title])
    for block in body.split("\n\n"):
        if block.strip():
            first_line = block.strip().splitlines()[0][:80]
            sheet.append([first_line, block.strip()])
    sheet.column_dimensions["A"].width = 28
    sheet.column_dimensions["B"].width = 90
    workbook.save(path)


def create_pptx_file(path: Path, title: str, body: str):
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("PowerPoint generation needs python-pptx. Install backend requirements again.") from exc

    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = "Generated by JARVIS"

    blocks = [block.strip() for block in body.split("\n\n") if block.strip()]
    for block in blocks[:8]:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        lines = block.splitlines()
        slide.shapes.title.text = lines[0][:80]
        slide.placeholders[1].text = "\n".join(lines[1:])[:900] or block[:900]

    presentation.save(path)


def create_txt_file(path: Path, body: str):
    path.write_text(body, encoding="utf-8")


def wants_active_app_typing(lowered: str):
    return any(phrase in lowered for phrase in [
        "type ",
        "paste ",
        "type it",
        "type this",
        "type that",
        "type what",
        "paste it",
        "paste this",
        "paste that",
        "write it there",
        "write that there",
        "put it there",
        "put that there",
        "fill this in",
        "fill that in",
        "into the app",
        "in the app",
        "in word",
        "in excel",
        "in powerpoint",
        "in notepad",
        "in vscode",
        "in vs code",
        "in the browser",
    ])


def wants_followup_app_typing(lowered: str, chat_id: str):
    if not latest_active_app(chat_id):
        return False
    return lowered.startswith(("write ", "type ", "paste ", "draft ", "fill ")) or any(phrase in lowered for phrase in [
        "write it",
        "put it",
        "add it",
        "enter it",
        "fill it",
        "there",
        "that app",
        "the app",
    ])


def extract_explicit_text_to_type(message: str):
    patterns = [
        r"(?is)\btype\s+['\"](.+?)['\"]\s*$",
        r"(?is)\bpaste\s+['\"](.+?)['\"]\s*$",
        r"(?is)\btype this\s*:\s*(.+)$",
        r"(?is)\bpaste this\s*:\s*(.+)$",
        r"(?is)^\s*type\s+(.+)$",
        r"(?is)^\s*paste\s+(.+)$",
    ]
    for pattern in patterns:
        match = re.search(pattern, message)
        if match:
            value = match.group(1).strip()
            if value and value.lower() not in {"it", "this", "that", "what"}:
                return value
    return None


def generate_action_draft(db: Session, chat_id: str, message: str):
    context = build_context(db, chat_id, message)
    if not context:
        return (
            "I do not have assignment or file context for this chat yet. "
            "Upload a file or add/read a page first, then ask me to type the required draft."
        )

    prompt = f"""You are JARVIS, a precise local study assistant.
Draft text that the user can paste into the currently focused app.
Use the saved assignment/file context. Do not mention that you are an AI.
Do not claim anything was submitted.
If the task is unclear, write a concise fill-in template with placeholders.

Saved context:
{context}

User request:
{message}

Paste-ready text:
"""
    try:
        response = requests.post(
            OLLAMA_URL,
            json={
                "model": DEFAULT_OLLAMA_MODEL,
                "prompt": prompt,
                "stream": False,
            },
            timeout=(10, OLLAMA_TIMEOUT_SECONDS),
        )
        response.raise_for_status()
        generated = response.json().get("response", "").strip()
        return generated or make_generation_text(message, context)
    except requests.RequestException:
        return make_generation_text(message, context)


def paste_text_into_active_window(text: str):
    if os.name != "nt":
        raise RuntimeError("Typing into the active app is currently implemented for Windows.")

    subprocess.run(
        [
            "powershell",
            "-NoProfile",
            "-Command",
            (
                "Add-Type -AssemblyName System.Windows.Forms; "
                "$text = [Console]::In.ReadToEnd(); "
                "Set-Clipboard -Value $text; "
                "Start-Sleep -Milliseconds 250; "
                "[System.Windows.Forms.SendKeys]::SendWait('^v')"
            ),
        ],
        input=text,
        text=True,
        check=True,
        capture_output=True,
    )


def mcq_session_key(chat_id: str):
    return str(chat_id)


def format_mcq_questions(questions):
    blocks = []
    for q_index, question in enumerate(questions, start=1):
        options = "\n".join(
            f"{chr(64 + option_index)}. {option.label}"
            for option_index, option in enumerate(question.options, start=1)
        )
        blocks.append(f"Q{q_index}. {question.question}\n{options}")
    return "\n\n".join(blocks)


def parse_mcq_answer_plan(raw_text: str, question_count: int):
    plan = {}
    json_match = re.search(r"\{.*\}", raw_text, flags=re.S)
    if json_match:
        try:
            data = json.loads(json_match.group(0))
            answers = data.get("answers", data if isinstance(data, list) else [])
            for item in answers:
                q_num = int(item.get("question", item.get("q", 0)))
                choice = str(item.get("choice", item.get("answer", ""))).strip().upper()
                if 1 <= q_num <= question_count and choice:
                    plan[q_num] = {
                        "choice": choice[0],
                        "reason": str(item.get("reason", "")).strip(),
                    }
        except Exception:
            pass

    if not plan:
        for match in re.finditer(r"(?i)q(?:uestion)?\s*(\d+)\D{0,20}(?:answer|choice)?\D{0,10}\b([A-Z])\b", raw_text):
            q_num = int(match.group(1))
            if 1 <= q_num <= question_count:
                plan[q_num] = {"choice": match.group(2).upper(), "reason": ""}

    return plan


def suggest_mcq_answers(questions, context: str):
    question_text = format_mcq_questions(questions)
    prompt = f"""You are JARVIS helping with visible multiple-choice questions.
Use the saved course/page context if helpful. Pick the best visible option for each question.
Return JSON only in this exact shape:
{{"answers":[{{"question":1,"choice":"A","reason":"short reason"}}]}}

Saved context:
{context or "No extra context."}

Visible questions and options:
{question_text}
"""
    try:
        response = requests.post(
            OLLAMA_URL,
            json={"model": DEFAULT_OLLAMA_MODEL, "prompt": prompt, "stream": False},
            timeout=(10, OLLAMA_TIMEOUT_SECONDS),
        )
        response.raise_for_status()
        raw_text = response.json().get("response", "").strip()
    except requests.RequestException:
        raw_text = ""

    plan = parse_mcq_answer_plan(raw_text, len(questions))
    lines = ["I found these visible MCQ questions and prepared a click plan:"]
    for q_index, question in enumerate(questions, start=1):
        answer = plan.get(q_index)
        if answer:
            lines.append(
                f"Q{q_index}: choose {answer['choice']}"
                + (f" - {answer['reason']}" if answer.get("reason") else "")
            )
        else:
            lines.append(f"Q{q_index}: I could not confidently choose an option.")

    lines.append("\nSay `mark these answers` and I will click the planned choices. I will not submit the quiz.")
    return plan, "\n".join(lines)


def start_mcq_browser_session(db: Session, chat_id: str, message: str):
    url_match = re.search(r"(https?://\S+|(?:www\.)?[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}(?:/\S*)?)", message)
    if not url_match:
        return True, "Send the quiz/page URL with the request, like `solve MCQ https://...`."

    key = mcq_session_key(chat_id)
    existing = ACTIVE_MCQ_SESSIONS.pop(key, None)
    if existing:
        try:
            existing["driver"].quit()
        except Exception:
            pass

    driver = create_driver(headless=False)
    url = url_match.group(1).rstrip(".,)")
    try:
        driver.get(url)
    except Exception as exc:
        driver.quit()
        return True, f"I could not open that MCQ page: {exc}"

    ACTIVE_MCQ_SESSIONS[key] = {
        "driver": driver,
        "questions": [],
        "plan": {},
        "url": url,
    }
    return True, (
        "I opened the quiz/page in a Selenium browser.\n\n"
        "Log in if needed and leave the questions visible, then say `read the MCQs now`."
    )


def read_mcq_browser_session(db: Session, chat_id: str):
    session = ACTIVE_MCQ_SESSIONS.get(mcq_session_key(chat_id))
    if not session:
        return True, "No active MCQ browser session. First say `solve MCQ <url>`."

    driver = session["driver"]
    questions = extract_mcq_questions(driver)
    if not questions:
        return True, "I could not find visible radio-button or checkbox MCQ options on the current page."

    context = build_context(db, chat_id, "multiple choice quiz answers")
    plan, response_text = suggest_mcq_answers(questions, context)
    session["questions"] = questions
    session["plan"] = plan
    return True, response_text


def mark_mcq_answers(chat_id: str):
    session = ACTIVE_MCQ_SESSIONS.get(mcq_session_key(chat_id))
    if not session or not session.get("questions"):
        return True, "No MCQ answer plan is ready. First say `read the MCQs now`."

    driver = session["driver"]
    questions = session["questions"]
    plan = session.get("plan", {})
    clicked = []
    skipped = []

    for q_num, answer in sorted(plan.items()):
        option_index = ord(answer["choice"].upper()[0]) - 65
        if q_num < 1 or q_num > len(questions) or option_index < 0 or option_index >= len(questions[q_num - 1].options):
            skipped.append(f"Q{q_num}")
            continue

        element_id = questions[q_num - 1].options[option_index].element_id
        try:
            element = next(
                input_el for input_el in driver.find_elements(By.CSS_SELECTOR, "input[type='radio'], input[type='checkbox']")
                if input_el.id == element_id
            )
            driver.execute_script("arguments[0].scrollIntoView({block: 'center'});", element)
            element.click()
            clicked.append(f"Q{q_num} {answer['choice'].upper()[0]}")
        except Exception:
            skipped.append(f"Q{q_num}")

    return True, (
        f"Marked: {', '.join(clicked) if clicked else 'none'}.\n"
        f"Skipped: {', '.join(skipped) if skipped else 'none'}.\n\n"
        "I did not submit the quiz. Review the selected answers before submitting."
    )


def close_mcq_browser_session(chat_id: str):
    session = ACTIVE_MCQ_SESSIONS.pop(mcq_session_key(chat_id), None)
    if not session:
        return True, "No active MCQ browser session is open."
    try:
        session["driver"].quit()
    except Exception:
        pass
    return True, "Closed the MCQ browser session."


def browser_session_key(chat_id: str):
    return str(chat_id)


def start_visible_browser_session(chat_id: str, url: str, label: str = "page"):
    key = browser_session_key(chat_id)
    existing = ACTIVE_BROWSER_SESSIONS.pop(key, None)
    if existing:
        try:
            existing["driver"].quit()
        except Exception:
            pass

    driver = create_driver(headless=False)
    try:
        driver.get(url)
    except Exception as exc:
        driver.quit()
        raise RuntimeError(f"I could not open {label} in Chrome: {exc}") from exc

    ACTIVE_BROWSER_SESSIONS[key] = {
        "driver": driver,
        "label": label,
        "url": url,
    }
    return (
        f"Opened {label} in a visible Chrome window.\n\n"
        "Use that window normally. When it is showing the page you want, say `read the open page` or `read what it says`."
    )


def collect_open_browser_page(driver):
    try:
        title = driver.title or "Open browser page"
    except Exception:
        title = "Open browser page"

    try:
        url = driver.current_url
    except Exception:
        url = ""

    try:
        body = driver.find_element(By.TAG_NAME, "body")
        text = body.text.strip()
    except Exception as exc:
        raise RuntimeError("I could not read the visible browser page. Make sure the Chrome window is still open.") from exc

    links = []
    pdf_links = []
    for link in driver.find_elements(By.CSS_SELECTOR, "a[href]"):
        href = link.get_attribute("href")
        if not href:
            continue
        label = link.text.strip() or link.get_attribute("aria-label") or href
        lower_href = href.lower()
        if ".pdf" in lower_href:
            link_type = "pdf"
            pdf_links.append(href)
        elif any(ext in lower_href for ext in [".docx", ".xlsx", ".pptx", ".csv", ".zip"]):
            link_type = "file"
        else:
            link_type = "web"
        links.append({"label": label[:160], "url": href, "type": link_type})

    seen = set()
    unique_links = []
    for link in links:
        if link["url"] in seen:
            continue
        seen.add(link["url"])
        unique_links.append(link)

    return {
        "title": title,
        "url": url,
        "text": text,
        "links": unique_links[:50],
        "pdf_links": sorted(set(pdf_links)),
    }


def read_visible_browser_session(db: Session, chat_id: str):
    session = ACTIVE_BROWSER_SESSIONS.get(browser_session_key(chat_id))
    if not session:
        return True, (
            "I do not have a live Chrome page to read yet. Say `open <site or search> in chrome`, navigate to the page you want, then say `read the open page`."
        )

    page = collect_open_browser_page(session["driver"])
    text = page["text"]
    if len(text.split()) < 10:
        return True, "I found very little readable text on the open Chrome page. Make sure the assignment content is visible, then ask me to read it again."

    source_label = page["title"] or session.get("label") or "Open browser page"
    source_text = f"URL: {page['url']}\n\n{text}"
    chunk_count = save_document_text(db, chat_id, source_label, source_text)
    combined_content = source_text
    imported_pdfs = []

    for pdf_url in page["pdf_links"][:5]:
        try:
            pdf_text = extract_pdf_url_text(pdf_url)
            pdf_chunks = save_document_text(db, chat_id, pdf_url, pdf_text)
            combined_content += "\n\n" + pdf_text
            imported_pdfs.append({"url": pdf_url, "chunks": pdf_chunks})
            chunk_count += pdf_chunks
        except Exception as exc:
            imported_pdfs.append({"url": pdf_url, "error": str(exc)})

    if page["links"]:
        link_summary = "\n\nLinked resources discovered on the open browser page:\n" + "\n".join(
            f"- [{link['type']}] {link['label']}: {link['url']}"
            for link in page["links"][:25]
        )
        chunk_count += save_document_text(db, chat_id, f"{source_label} links", link_summary)
        combined_content += link_summary

    return True, attachment_reply(
        source_label,
        chunk_count,
        combined_content,
        source_kind="open Chrome page",
        links=page["links"],
        imported_pdfs=imported_pdfs,
    )


def close_visible_browser_session(chat_id: str):
    session = ACTIVE_BROWSER_SESSIONS.pop(browser_session_key(chat_id), None)
    if not session:
        return True, "No live Chrome browser session is open."
    try:
        session["driver"].quit()
    except Exception:
        pass
    return True, "Closed the live Chrome browser session."


def wants_live_browser_open(lowered: str):
    return lowered.startswith("open ") and any(word in lowered for word in ["chrome", "browser", "webpage", "website", "page"])


def wants_open_browser_read(lowered: str):
    read_terms = ["read", "see what it says", "what does it say", "what it says", "scan", "summarize", "look at"]
    browser_terms = ["chrome", "browser", "webpage", "website", "open page", "current page", "assignment"]
    return any(term in lowered for term in read_terms) and any(term in lowered for term in browser_terms)


def run_desktop_action(db: Session, chat_id: str, message: str):
    text = message.strip()
    lowered = text.lower()

    url_match = re.search(r"(https?://\S+|(?:www\.)?[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}(?:/\S*)?)", text)
    app_label, app_url = extract_known_web_app(text)
    browser_target = extract_browser_open_target(text)
    if (
        wants_live_browser_open(lowered)
        or (lowered.startswith("open ") and app_url)
        or (lowered.startswith("open ") and url_match)
        or (lowered.startswith("open ") and browser_target and looks_like_web_target(browser_target))
    ):
        if url_match:
            target_url = url_match.group(1).rstrip(".,)")
            label = app_label.title() if app_label else "that page"
        elif app_url:
            target_url = app_url
            label = app_label.title()
        elif browser_target and browser_target.lower() not in {"chrome", "browser"}:
            target_url = browser_url_for_target(browser_target)
            label = browser_target
        else:
            target_url = None
            label = None

        if target_url:
            if not re.match(r"^https?://", target_url, re.I):
                target_url = f"https://{target_url}"
            return True, start_visible_browser_session(chat_id, target_url, label)

    if ACTIVE_BROWSER_SESSIONS.get(browser_session_key(chat_id)) and any(
        term in lowered for term in ["read", "see what it says", "what does it say", "what it says", "scan", "summarize", "look at"]
    ):
        return read_visible_browser_session(db, chat_id)

    if wants_camera_capture(lowered, chat_id):
        return capture_camera_photo(chat_id)

    if wants_calculator_action(lowered, chat_id):
        return run_calculator_action(chat_id, text)

    if wants_generic_key_action(lowered, chat_id):
        return run_generic_key_action(chat_id, text)

    if (wants_active_app_typing(lowered) or wants_followup_app_typing(lowered, chat_id)) and not wants_open_browser_read(lowered):
        explicit_text = extract_explicit_text_to_type(text)
        draft = explicit_text or generate_action_draft(db, chat_id, text)
        paste_text_into_active_window(draft)
        active = latest_active_app(chat_id)
        target_label = active["label"].title() if active else "the currently focused app/window"
        return True, (
            f"I pasted the draft into {target_label}.\n\n"
            "Review it before saving, submitting, or sending."
        )

    if wants_open_browser_read(lowered):
        return read_visible_browser_session(db, chat_id)

    if any(phrase in lowered for phrase in ["close live browser", "close chrome session", "close open browser", "close canvas browser"]):
        return close_visible_browser_session(chat_id)

    if url_match and (lowered.startswith("open ") or any(phrase in lowered for phrase in ["open website", "open link", "open page", "go to"])):
        url = open_website(url_match.group(1).rstrip(".,)"))
        return True, f"Opened this page in your browser:\n{url}"

    search_query = extract_search_query(text)
    if search_query:
        url = browser_url_for_target(search_query)
        return True, start_visible_browser_session(chat_id, url, f"search for {search_query}")

    if any(phrase in lowered for phrase in ["solve mcq", "solve quiz", "mcq page", "quiz page", "answer mcq"]):
        return start_mcq_browser_session(db, chat_id, text)

    if any(phrase in lowered for phrase in ["read the mcqs", "read mcqs", "read the quiz", "read quiz questions"]):
        return read_mcq_browser_session(db, chat_id)

    if any(phrase in lowered for phrase in ["mark these answers", "mark right answers", "click these answers", "select these answers"]):
        return mark_mcq_answers(chat_id)

    if any(phrase in lowered for phrase in ["close quiz browser", "close mcq browser", "close selenium quiz"]):
        return close_mcq_browser_session(chat_id)

    if wants_active_app_typing(lowered):
        explicit_text = extract_explicit_text_to_type(text)
        draft = explicit_text or generate_action_draft(db, chat_id, text)
        paste_text_into_active_window(draft)
        return True, (
            "I pasted the draft into the currently focused app/window.\n\n"
            "Review it before saving, submitting, or sending."
        )

    if lowered.startswith("open ") or " open " in lowered:
        app_label, app_command = allowed_app_command(lowered)
        if app_command:
            launch_allowed_app(app_command)
            remember_active_app(chat_id, app_label, app_command)
            return True, f"Opened {app_label.title()}."

        requested_app = extract_open_app_name(text)
        if requested_app:
            try:
                launch_named_app(requested_app)
                remember_active_app(chat_id, requested_app, {"command": requested_app})
                return True, f"Asked Windows to open {requested_app}."
            except RuntimeError:
                url = browser_url_for_target(requested_app)
                return True, start_visible_browser_session(chat_id, url, requested_app)

    wants_file = any(phrase in lowered for phrase in [
        "make a word",
        "create a word",
        "write in word",
        "make a doc",
        "create a doc",
        "make an excel",
        "create an excel",
        "spreadsheet",
        "make a powerpoint",
        "create a powerpoint",
        "make slides",
        "create slides",
        "make a text file",
        "create a text file",
    ])

    if wants_file:
        context = build_context(db, chat_id, text)
        title = "JARVIS Draft"
        body = make_generation_text(text, context)
        GENERATED_DIR.mkdir(parents=True, exist_ok=True)

        if any(word in lowered for word in ["excel", "spreadsheet", "xlsx"]):
            suffix = ".xlsx"
            creator = create_xlsx_file
        elif any(word in lowered for word in ["powerpoint", "power point", "slides", "pptx"]):
            suffix = ".pptx"
            creator = create_pptx_file
        elif any(word in lowered for word in ["text file", "txt"]):
            suffix = ".txt"
            creator = create_txt_file
        else:
            suffix = ".docx"
            creator = create_docx_file

        filename = f"{safe_filename(text)}_{uuid4().hex[:8]}{suffix}"
        path = GENERATED_DIR / filename
        if suffix == ".txt":
            creator(path, body)
        else:
            creator(path, title, body)

        if any(phrase in lowered for phrase in ["open it", "open the file", "open word", "open excel", "open powerpoint"]):
            open_local_target(str(path))

        return True, (
            f"I created this file:\n{path}\n\n"
            f"Download link: {generated_file_url(path)}\n\n"
            "It includes the request plus the most relevant saved assignment/context I found. "
            "Ask me to revise it, make it more formal, turn it into bullets, or add more detail."
        )

    if any(phrase in lowered for phrase in ["submit", "apply for", "fill job application", "fill application", "click answer", "mark answer"]):
        return True, (
            "I can help prepare answers and guide the browser workflow, but I will not submit, apply, or click final answers automatically yet. "
            "For that we need a confirmation-based browser action mode so every important click is approved by you first."
        )

    return False, ""


def classify_assignment(text: str):
    lowered = text.lower()
    explicit_mcq_signals = ["multiple choice", "choose the correct", "select one", "select all"]
    option_lines = len(re.findall(r"(?im)^\s*[a-d][\).]\s+\S+", text))
    coding_signals = ["write a program", "function", "class ", "python", "java", "javascript", "sql", "algorithm", "debug"]
    writing_signals = ["essay", "paragraph", "discussion post", "reflection", "write about", "cite", "mla", "apa"]

    if any(signal in lowered for signal in explicit_mcq_signals) or option_lines >= 4:
        return "MCQ"
    if any(signal in lowered for signal in coding_signals):
        return "Coding"
    if any(signal in lowered for signal in writing_signals):
        return "Writing"
    return "General Assignment"


def describe_source_signals(content: str, label: str = ""):
    lowered = content.lower()
    signals = []
    option_lines = len(re.findall(r"(?im)^\s*[a-d][\).]\s+\S+", content))

    if any(phrase in lowered for phrase in ["rubric", "grading criteria", "grade criteria", "points possible", "total points"]):
        signals.append("rubric or grading details")
    if any(phrase in lowered for phrase in ["due date", "deadline", "available until", "submit by", "submission due"]):
        signals.append("deadline or submission details")
    if any(word in lowered for word in ["multiple choice", "choose the correct", "select one", "select all"]) or option_lines >= 4:
        signals.append("multiple-choice questions")
    if any(phrase in lowered for phrase in [
        "write a program",
        "write code",
        "implement a",
        "create a function",
        "debug the",
        "programming assignment",
        "coding assignment",
        "algorithm assignment",
    ]):
        signals.append("coding requirements")
    if any(phrase in lowered for phrase in [
        "write an essay",
        "write a paragraph",
        "discussion post",
        "reflection paper",
        "cite your sources",
        "mla format",
        "apa format",
    ]):
        signals.append("writing requirements")
    if any(phrase in lowered for phrase in ["template file", "starter file", "sample template", "use the template", "provided template"]):
        signals.append("templates or examples")

    return signals


def extract_content_snapshot(content: str):
    lines = []
    seen = set()
    for raw_line in content.splitlines():
        line = re.sub(r"\s+", " ", raw_line).strip()
        if not line or len(line) < 3 or len(line) > 120:
            continue
        key = line.lower()
        if key in seen:
            continue
        seen.add(key)
        lines.append(line)
        if len(lines) >= 5:
            break

    dates = sorted(set(re.findall(
        r"\b(?:\d{1,2}[/-]\d{1,2}[/-]\d{2,4}|(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\.?\s+\d{1,2},?\s+\d{2,4})\b",
        content,
        flags=re.I,
    )))[:5]
    emails = sorted(set(re.findall(r"\b[\w.+-]+@[\w.-]+\.\w+\b", content)))[:5]
    urls = sorted(set(re.findall(r"https?://\S+", content)))[:5]

    snapshot = []
    if lines:
        snapshot.append("Visible lines:\n" + "\n".join(f"- {line}" for line in lines))
    if dates:
        snapshot.append("Dates found: " + ", ".join(dates))
    if emails:
        snapshot.append("Emails found: " + ", ".join(emails))
    if urls:
        snapshot.append("URLs found: " + ", ".join(urls))
    return snapshot


def make_preview(content: str, limit=500):
    preview = re.sub(r"\s+", " ", content).strip()
    if len(preview) <= limit:
        return preview
    return f"{preview[:limit].rstrip()}..."


def attachment_reply(label: str, chunk_count: int, content: str, source_kind: str = "file", links=None, imported_pdfs=None):
    if is_image_filename(label):
        readable = content.strip()
        if len(readable) < 12:
            return (
                f"I added {label} to this chat, but OCR found very little readable text: "
                f"`{readable or 'nothing'}`.\n\n"
                "For visual understanding, choose the Vision model from the model dropdown and ask what is in the image."
            )

        return (
            f"I added {label} to this chat and extracted readable image text with OCR.\n\n"
            f"Extracted text preview: {readable[:300]}\n\n"
            "Ask me about the text I found, or choose the Vision model to ask about the image itself."
        )

    word_count = len(content.split())
    signals = describe_source_signals(content, label)
    snapshot = extract_content_snapshot(content)
    lines = [
        f"I read {label} and saved {chunk_count} searchable chunks.",
        f"Source: {source_kind}. Words found: {word_count}.",
    ]

    if signals:
        lines.append(f"Strong signals found: {', '.join(signals)}.")

    if snapshot:
        lines.append("\n\n".join(snapshot))

    if links:
        pdf_count = len([link for link in links if link.get("type") == "pdf"])
        file_count = len([link for link in links if link.get("type") == "file"])
        web_count = len([link for link in links if link.get("type") == "web"])
        pieces = []
        if pdf_count:
            pieces.append(f"{pdf_count} PDF link{'s' if pdf_count != 1 else ''}")
        if file_count:
            pieces.append(f"{file_count} file link{'s' if file_count != 1 else ''}")
        if web_count:
            pieces.append(f"{web_count} web link{'s' if web_count != 1 else ''}")
        if pieces:
            lines.append(f"Linked resources found: {', '.join(pieces)}.")

    if imported_pdfs:
        successful = [item for item in imported_pdfs if "chunks" in item]
        failed = [item for item in imported_pdfs if "error" in item]
        if successful:
            lines.append(f"I also imported {len(successful)} linked PDF{'s' if len(successful) != 1 else ''}.")
        if failed:
            lines.append(f"{len(failed)} linked PDF{'s' if len(failed) != 1 else ''} could not be imported.")

    preview = make_preview(content)
    if preview:
        lines.append(f"Preview: {preview}")

    lines.append("Ask me what you want to do with it: summarize, answer questions, draft from it, extract requirements, make a file, or compare it with another source.")
    return "\n\n".join(lines)


def vision_attachment_reply(filename: str, vision_text: str, ocr_text: str):
    readable = ocr_text.strip()
    if not readable or "No readable OCR text" in readable:
        ocr_summary = "OCR did not find readable text."
    elif len(readable) < 12:
        ocr_summary = f"OCR found very little readable text: `{readable}`."
    else:
        ocr_summary = f"OCR text found: {readable[:500]}"

    lines = [
        f"I added {filename} to this chat and inspected it with the selected vision model.",
        ocr_summary,
        f"Vision read: {vision_text.strip()}",
    ]

    lines.append("I saved the OCR text as chat context, and while a vision model is selected I can also inspect the image directly for follow-up questions.")
    return "\n\n".join(lines)


def analyze_uploaded_image_with_ollama(filename: str, file_bytes: bytes, model_name: str):
    prompt = (
        "Describe this uploaded image clearly and concisely. "
        "If it contains text, transcribe the visible text. "
        "If it appears to be a signature, logo, document, screenshot, or photo, say that directly."
    )
    response = requests.post(
        OLLAMA_URL,
        json={
            "model": model_name,
            "prompt": prompt,
            "images": [base64.b64encode(file_bytes).decode("utf-8")],
            "stream": False,
        },
        timeout=(10, OLLAMA_TIMEOUT_SECONDS),
    )
    response.raise_for_status()
    return response.json().get("response", "").strip()


@app.get("/health")
def health():
    return {"status": "ok"}


@app.get("/models")
def models():
    try:
        installed_models = fetch_installed_ollama_models()
        model_names = [model["name"] for model in installed_models]
    except requests.RequestException:
        model_names = sorted(ALLOWED_OLLAMA_MODELS)

    return {
        "default": DEFAULT_OLLAMA_MODEL,
        "models": model_names,
    }


@app.get("/capabilities")
def capabilities():
    return {
        "chat": ["streaming local Ollama answers", "saved chat history", "document memory"],
        "files": ["PDF", "DOCX", "XLSX", "TXT/CSV/MD", "image OCR", "vision model support"],
        "web": ["Selenium page reading", "manual login mode", "linked PDF import", "MCQ extraction"],
        "actions": [
            "open installed apps",
            "open websites",
            "web search",
            "create Word/Excel/PowerPoint/TXT files",
            "paste generated drafts into the focused app",
            "controlled MCQ marking without final submit",
        ],
        "safety": ["does not submit quizzes/forms automatically", "asks for user review before final actions"],
    }


@app.get("/generated-files")
def generated_files(user=Depends(get_current_user)):
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    files = []
    for path in sorted(GENERATED_DIR.glob("*"), key=lambda item: item.stat().st_mtime, reverse=True):
        if path.is_file():
            stat = path.stat()
            files.append({
                "name": path.name,
                "size": stat.st_size,
                "url": generated_file_url(path),
            })
    return files[:50]


@app.get("/generated/{filename}")
def download_generated_file(filename: str, user=Depends(get_current_user_optional)):
    path = resolve_generated_file(filename)
    return FileResponse(path, filename=path.name)


def normalize_model_name(model_name: str | None):
    selected = model_name or DEFAULT_OLLAMA_MODEL
    return MODEL_ALIASES.get(selected, selected)


def is_vision_model(model_name: str):
    lowered = model_name.lower()
    return model_name in VISION_OLLAMA_MODELS or any(
        keyword in lowered for keyword in VISION_MODEL_KEYWORDS
    )


def fetch_installed_ollama_models(timeout=3):
    response = requests.get(f"{OLLAMA_BASE_URL}/api/tags", timeout=timeout)
    response.raise_for_status()
    models = []
    for model in response.json().get("models", []):
        name = model.get("name")
        if not name:
            continue
        models.append({
            "name": name,
            "size": model.get("size"),
            "modified_at": model.get("modified_at"),
            "details": model.get("details") or {},
        })
    return sorted(models, key=lambda item: item["name"].lower())


def model_status_entry(model):
    name = model["name"]
    return {
        "available": True,
        "status": "active",
        "vision": is_vision_model(name),
        "size": model.get("size"),
        "modified_at": model.get("modified_at"),
        "family": (model.get("details") or {}).get("family"),
    }


@app.get("/models/status")
def model_status():
    try:
        installed_models = fetch_installed_ollama_models()
    except requests.RequestException:
        return {
            "ollama": "offline",
            "default": DEFAULT_OLLAMA_MODEL,
            "models": {
                model_name: {"available": False, "status": "ollama_offline"}
                for model_name in sorted(ALLOWED_OLLAMA_MODELS)
            },
        }

    statuses = {
        model["name"]: model_status_entry(model)
        for model in installed_models
    }
    if DEFAULT_OLLAMA_MODEL not in statuses and statuses:
        default_model = next(
            (name for name in statuses if name.startswith(DEFAULT_OLLAMA_MODEL.split(":")[0])),
            next(iter(statuses)),
        )
    else:
        default_model = DEFAULT_OLLAMA_MODEL

    return {
        "ollama": "online",
        "default": default_model,
        "models": statuses,
    }


@app.post("/signup")
def signup(user: AuthRequest, db: Session = Depends(get_db)):
    email = user.email.lower()
    existing = db.query(User).filter(User.email == email).first()
    if existing:
        raise HTTPException(status_code=409, detail="User already exists")

    new_user = User(id=str(uuid4()), email=email, password=hash_password(user.password))
    db.add(new_user)
    db.commit()
    return {"msg": "user created"}


@app.post("/login")
def login(user: AuthRequest, db: Session = Depends(get_db)):
    email = user.email.lower()
    db_user = db.query(User).filter(User.email == email).first()
    if not db_user or not db_user.password or not verify_password(user.password, db_user.password):
        raise HTTPException(status_code=401, detail="Invalid credentials")

    return {
        "access_token": create_access_token({"user": db_user.email}),
        "refresh_token": create_refresh_token({"user": db_user.email}),
    }


@app.post("/forgot-password")
def forgot_password(data: ForgotPasswordRequest, db: Session = Depends(get_db)):
    email = data.email.lower().strip()
    db_user = db.query(User).filter(User.email == email).first()
    if not db_user:
        raise HTTPException(status_code=404, detail="Account does not exist")
    return {
        "msg": "Account found",
        "email": db_user.email,
        "masked_email": mask_email(db_user.email),
    }


@app.post("/send-reset-code")
def send_reset_code(data: ForgotPasswordRequest, db: Session = Depends(get_db)):
    email = data.email.lower().strip()
    db_user = db.query(User).filter(User.email == email).first()
    if not db_user:
        raise HTTPException(status_code=404, detail="Account does not exist")

    code = create_reset_code(db_user.email)
    try:
        send_reset_code_email(db_user.email, code)
    except RuntimeError as exc:
        PASSWORD_RESET_CODES.pop(db_user.email, None)
        raise HTTPException(status_code=503, detail=str(exc)) from exc
    except Exception as exc:
        PASSWORD_RESET_CODES.pop(db_user.email, None)
        raise HTTPException(status_code=502, detail="Could not send reset email. Check Resend settings.") from exc

    return {
        "msg": "Email reset code has been sent to your email",
        "masked_email": mask_email(db_user.email),
        "expires_minutes": RESET_CODE_EXPIRE_MINUTES,
    }


@app.post("/verify-reset-code")
def verify_reset_code(data: ResetCodeRequest, db: Session = Depends(get_db)):
    email = data.email.lower().strip()
    reset = PASSWORD_RESET_CODES.get(email)
    if not reset or reset["expires_at"] < datetime.utcnow():
        PASSWORD_RESET_CODES.pop(email, None)
        raise HTTPException(status_code=400, detail="Reset code is invalid or expired")

    reset["attempts"] += 1
    if reset["attempts"] > 5:
        PASSWORD_RESET_CODES.pop(email, None)
        raise HTTPException(status_code=429, detail="Too many attempts. Request a new code.")

    if reset["code"] != data.code.strip():
        raise HTTPException(status_code=400, detail="Verification code is wrong")

    db_user = db.query(User).filter(User.email == email).first()
    if not db_user:
        raise HTTPException(status_code=404, detail="Account not found")

    token = secrets.token_urlsafe(32)
    PASSWORD_RESET_TOKENS[token] = {
        "email": email,
        "expires_at": datetime.utcnow() + timedelta(minutes=PASSWORD_RESET_EXPIRE_MINUTES),
    }
    reset["verified"] = True
    return {"msg": "ok", "reset_token": token}


@app.post("/reset-password")
def reset_password(data: ResetPasswordRequest, db: Session = Depends(get_db)):
    reset = PASSWORD_RESET_TOKENS.get(data.token)
    if not reset or reset["expires_at"] < datetime.utcnow():
        PASSWORD_RESET_TOKENS.pop(data.token, None)
        raise HTTPException(status_code=400, detail="Reset token is invalid or expired")

    password = data.password.strip()
    if len(password) < 6:
        raise HTTPException(status_code=422, detail="Password must be at least 6 characters")

    db_user = db.query(User).filter(User.email == reset["email"]).first()
    if not db_user:
        PASSWORD_RESET_TOKENS.pop(data.token, None)
        raise HTTPException(status_code=404, detail="Account not found")

    db_user.password = hash_password(password)
    db.commit()
    PASSWORD_RESET_TOKENS.pop(data.token, None)
    return {"msg": "Password updated"}


@app.post("/change-password")
def change_password(data: ChangePasswordRequest, db: Session = Depends(get_db)):
    reset = PASSWORD_RESET_TOKENS.get(data.reset_token)
    email = data.email.lower().strip()
    if not reset or reset["email"] != email or reset["expires_at"] < datetime.utcnow():
        PASSWORD_RESET_TOKENS.pop(data.reset_token, None)
        raise HTTPException(status_code=400, detail="Password reset session is invalid or expired")

    password = data.password.strip()
    if len(password) < 6:
        raise HTTPException(status_code=422, detail="Password must be at least 6 characters")

    db_user = db.query(User).filter(User.email == email).first()
    if not db_user:
        raise HTTPException(status_code=404, detail="Account not found")

    db_user.password = hash_password(password)
    db.commit()
    PASSWORD_RESET_TOKENS.pop(data.reset_token, None)
    PASSWORD_RESET_CODES.pop(email, None)
    return {"msg": "ok"}


@app.get("/oauth/{provider}/start")
def oauth_start(provider: str):
    provider = provider.lower()
    config = OAUTH_PROVIDERS.get(provider)
    if not config:
        raise HTTPException(status_code=404, detail="Unsupported login provider")

    client_id = os.getenv(config["client_id_env"])
    if not client_id or not os.getenv(config["client_secret_env"]):
        raise HTTPException(status_code=503, detail=f"{provider.title()} login is not configured")

    state = secrets.token_urlsafe(24)
    OAUTH_STATES[state] = {
        "provider": provider,
        "expires_at": datetime.utcnow() + timedelta(minutes=OAUTH_STATE_EXPIRE_MINUTES),
    }
    params = {
        "client_id": client_id,
        "redirect_uri": oauth_redirect_uri(provider),
        "response_type": "code",
        "scope": config["scope"],
        "state": state,
    }
    if provider == "google":
        params["access_type"] = "offline"
        params["prompt"] = "select_account"
    return RedirectResponse(f"{config['authorize_url']}?{urlencode(params)}")


@app.get("/oauth/{provider}/callback")
def oauth_callback(provider: str, code: str | None = None, state: str | None = None, db: Session = Depends(get_db)):
    provider = provider.lower()
    config = OAUTH_PROVIDERS.get(provider)
    state_data = OAUTH_STATES.pop(state or "", None)
    if not config or not code or not state_data or state_data["provider"] != provider:
        raise HTTPException(status_code=400, detail="Invalid login response")
    if state_data["expires_at"] < datetime.utcnow():
        raise HTTPException(status_code=400, detail="Login session expired")

    token_res = requests.post(
        config["token_url"],
        data={
            "client_id": os.getenv(config["client_id_env"]),
            "client_secret": os.getenv(config["client_secret_env"]),
            "code": code,
            "redirect_uri": oauth_redirect_uri(provider),
            "grant_type": "authorization_code",
        },
        headers={"Accept": "application/json"},
        timeout=20,
    )
    token_res.raise_for_status()
    token_data = token_res.json()
    access_token = token_data.get("access_token")
    if not access_token:
        raise HTTPException(status_code=400, detail="Provider did not return an access token")

    user_res = requests.get(
        config["userinfo_url"],
        headers={"Authorization": f"Bearer {access_token}", "Accept": "application/json"},
        timeout=20,
    )
    user_res.raise_for_status()
    profile = user_res.json()
    email = profile.get("email")

    if not email:
        raise HTTPException(status_code=400, detail="Provider account did not include a verified email")

    db_user = get_or_create_oauth_user(db, email.lower())
    app_access = create_access_token({"user": db_user.email})
    app_refresh = create_refresh_token({"user": db_user.email})
    return oauth_success_page(app_access, app_refresh)


@app.post("/refresh")
def refresh(credentials: HTTPAuthorizationCredentials = Depends(security)):
    data = decode_token(credentials.credentials)
    if not data or data.get("type") != "refresh":
        raise HTTPException(status_code=401, detail="Invalid refresh token")
    return {"access_token": create_access_token({"user": data["user"]})}


@app.post("/chat/create")
def create_chat(user=Depends(get_current_user), db: Session = Depends(get_db)):
    chat = Chat(id=str(uuid4()), title="New Chat", user_id=user)
    db.add(chat)
    db.commit()
    db.refresh(chat)
    return chat


@app.get("/chats")
def get_chats(user=Depends(get_current_user), db: Session = Depends(get_db)):
    return db.query(Chat).filter(Chat.user_id == user).all()


@app.get("/history/{chat_id}")
def history(chat_id: str, user=Depends(get_current_user), db: Session = Depends(get_db)):
    get_owned_chat(chat_id, user, db)
    return db.query(Message).filter(Message.chat_id == chat_id).all()


@app.put("/chat/{chat_id}")
def rename_chat(
    chat_id: str,
    data: RenameRequest,
    user=Depends(get_current_user),
    db: Session = Depends(get_db),
):
    chat = get_owned_chat(chat_id, user, db)
    chat.title = data.title.strip() or chat.title
    db.commit()
    return {"msg": "updated"}


@app.delete("/chat/{chat_id}")
def delete_chat(chat_id: str, user=Depends(get_current_user), db: Session = Depends(get_db)):
    chat = get_owned_chat(chat_id, user, db)
    db.query(Message).filter(Message.chat_id == chat_id).delete()
    db.query(DocumentChunk).filter(DocumentChunk.chat_id == chat_id).delete()
    db.query(ImageAttachment).filter(ImageAttachment.chat_id == chat_id).delete()
    db.delete(chat)
    db.commit()
    return {"msg": "deleted"}


@app.post("/upload/{chat_id}")
def upload(
    chat_id: str,
    file: UploadFile = File(...),
    model: str | None = Form(None),
    user=Depends(get_current_user),
    db: Session = Depends(get_db),
):
    get_owned_chat(chat_id, user, db)
    file_bytes = file.file.read()
    is_image = is_image_filename(file.filename)
    if is_image:
        save_image_attachment(db, chat_id, file.filename, file_bytes)

    content = extract_file_text(file.filename, file_bytes).strip()

    if not content and not is_image:
        raise HTTPException(status_code=422, detail="No text extracted from file")
    if not content and is_image:
        content = "No readable OCR text was extracted from this image."

    chunks = split_text(content)
    save_document_text(db, chat_id, file.filename, content)
    assistant_content = attachment_reply(file.filename, len(chunks), content, source_kind="uploaded file")
    selected_model = normalize_model_name(model)

    if is_image and is_vision_model(selected_model):
        try:
            installed_model_names = {item["name"] for item in fetch_installed_ollama_models()}
            if selected_model in installed_model_names:
                vision_text = analyze_uploaded_image_with_ollama(file.filename, file_bytes, selected_model)
                if vision_text:
                    assistant_content = vision_attachment_reply(file.filename, vision_text, content)
            else:
                assistant_content += (
                    f"\n\n`{selected_model}` is selected but is not installed in Ollama. "
                    "Choose an installed vision model from the dropdown."
                )
        except requests.Timeout:
            assistant_content += (
                f"\n\nI tried to inspect it with `{selected_model}`, but the vision model took too long to respond. "
                "The OCR text is still saved, so you can ask about the image text or try the vision model again."
            )
        except requests.RequestException:
            assistant_content += (
                f"\n\nI tried to inspect it with `{selected_model}`, but Ollama did not return a vision response. "
                "Make sure Ollama is running, then try again."
            )

    db.add_all([
        Message(
            id=str(uuid4()),
            chat_id=chat_id,
            role="user",
            content=f"Uploaded file: {file.filename}",
        ),
        Message(
            id=str(uuid4()),
            chat_id=chat_id,
            role="assistant",
            content=assistant_content,
        ),
    ])
    db.commit()

    return {
        "msg": "uploaded",
        "filename": file.filename,
        "chunks": len(chunks),
        "assistant_message": assistant_content,
        "source_signals": describe_source_signals(content, file.filename),
        "preview": content[:500],
    }


@app.post("/scrape/{chat_id}")
def scrape_assignment(
    chat_id: str,
    data: ScrapeRequest,
    user=Depends(get_current_user),
    db: Session = Depends(get_db),
):
    get_owned_chat(chat_id, user, db)
    try:
        page = scrape_assignment_page(
            data.url,
            headless=data.headless,
            wait_seconds=max(0, min(data.wait_seconds, 20)),
            manual_login=data.manual_login,
            login_wait_seconds=max(10, min(data.login_wait_seconds, 300)),
        )
    except RuntimeError as exc:
        raise HTTPException(status_code=422, detail=str(exc)) from exc

    chunk_count = save_document_text(db, chat_id, page.title, page.text)
    combined_content = page.text
    imported_pdfs = []

    if data.include_pdfs:
        for pdf_url in page.pdf_links[:5]:
            try:
                pdf_text = extract_pdf_url_text(pdf_url)
                pdf_chunks = save_document_text(db, chat_id, pdf_url, pdf_text)
                combined_content += "\n\n" + pdf_text
                imported_pdfs.append({"url": pdf_url, "chunks": pdf_chunks})
                chunk_count += pdf_chunks
            except Exception as exc:
                imported_pdfs.append({"url": pdf_url, "error": str(exc)})

    if page.links:
        link_summary = "\n\nLinked resources discovered on the page:\n" + "\n".join(
            f"- [{link['type']}] {link['label']}: {link['url']}"
            for link in page.links[:25]
        )
        chunk_count += save_document_text(db, chat_id, f"{page.title} links", link_summary)
        combined_content += link_summary

    assistant_content = attachment_reply(
        page.title,
        chunk_count,
        combined_content,
        source_kind="webpage",
        links=page.links,
        imported_pdfs=imported_pdfs,
    )
    db.add_all([
        Message(
            id=str(uuid4()),
            chat_id=chat_id,
            role="user",
            content=f"Added page: {data.url}",
        ),
        Message(
            id=str(uuid4()),
            chat_id=chat_id,
            role="assistant",
            content=assistant_content,
        ),
    ])
    db.commit()
    return {
        "msg": "scraped",
        "title": page.title,
        "chunks": chunk_count,
        "assistant_message": assistant_content,
        "source_signals": describe_source_signals(combined_content, page.title),
        "pdf_links": page.pdf_links,
        "links": page.links,
        "imported_pdfs": imported_pdfs,
        "preview": page.text[:500],
    }


@app.post("/action")
def assistant_action(
    data: ActionRequest,
    user=Depends(get_current_user),
    db: Session = Depends(get_db),
):
    get_owned_chat(data.chat_id, user, db)
    try:
        handled, assistant_content = run_desktop_action(db, data.chat_id, data.message)
    except RuntimeError as exc:
        handled = True
        assistant_content = str(exc)
    except Exception as exc:
        handled = True
        assistant_content = f"I could not complete that desktop action: {exc}"

    if handled:
        db.add_all([
            Message(id=str(uuid4()), chat_id=data.chat_id, role="user", content=data.message),
            Message(id=str(uuid4()), chat_id=data.chat_id, role="assistant", content=assistant_content),
        ])
        db.commit()

    return {
        "handled": handled,
        "assistant_message": assistant_content,
    }


@app.post("/stream")
def stream(
    data: StreamRequest,
    user=Depends(get_current_user),
    db: Session = Depends(get_db),
):
    get_owned_chat(data.chat_id, user, db)
    selected_model = normalize_model_name(data.model)
    try:
        installed_model_names = {model["name"] for model in fetch_installed_ollama_models()}
    except requests.RequestException as exc:
        raise HTTPException(status_code=503, detail="Ollama is offline. Start Ollama, then try again.") from exc

    if selected_model not in installed_model_names:
        raise HTTPException(status_code=400, detail="Choose an installed Ollama model from the dropdown")

    def generate():
        full_response = ""
        previous_messages = db.query(Message).filter(Message.chat_id == data.chat_id).all()
        history_text = "\n".join(
            f"{msg.role}: {msg.content}" for msg in previous_messages[-8:]
        )
        context = build_context(db, data.chat_id, data.message)
        image_payloads = get_latest_image_payloads(db, data.chat_id) if is_vision_model(selected_model) else []

        source_signals = describe_source_signals(context)

        prompt = f"""You are JARVIS, a calm, precise AI study assistant.
Be proactive, concise, and useful. Use uploaded or scraped assignment context when it is relevant.
If the context does not contain the answer, say what is missing and ask for the needed detail.
If a vision image is provided, inspect it directly and answer about what you can see.
If no vision image is provided, never invent details about an uploaded image; you only know the OCR/extracted text shown in context.
If a user asks what an image is and you do not receive image data, say to switch to the Vision model.
For MCQs, explain the reasoning and point to the best option from the provided choices.
For coding, break down the requirements and help implement or debug step by step.
For writing, help outline, draft, revise, and cite responsibly.
Do not claim you submitted anything or accessed hidden answers.

Uploaded document context:
{context or "No uploaded document context for this chat."}

Source signals noticed:
{", ".join(source_signals) if source_signals else "No special assignment pattern detected. Use the context directly."}

Recent conversation:
{history_text or "No previous messages."}

User question:
{data.message}
"""

        try:
            ollama_payload = {
                "model": selected_model,
                "prompt": prompt,
                "stream": True,
            }
            if image_payloads:
                ollama_payload["images"] = image_payloads

            response = requests.post(
                OLLAMA_URL,
                json=ollama_payload,
                stream=True,
                timeout=(10, OLLAMA_TIMEOUT_SECONDS),
            )
            response.raise_for_status()

            for line in response.iter_lines():
                if not line:
                    continue
                chunk = json.loads(line.decode("utf-8"))
                token = chunk.get("response", "")
                full_response += token
                yield token
        except requests.HTTPError as exc:
            status_code = exc.response.status_code if exc.response is not None else None
            if status_code == 404:
                full_response = (
                    f"The selected model `{selected_model}` is not installed in Ollama. "
                    f"Run `ollama pull {selected_model}` or choose a model marked Active."
                )
            else:
                full_response = (
                    f"Ollama returned an error for `{selected_model}`. "
                    "Try another active model or restart Ollama."
                )
            yield full_response
        except requests.ConnectionError:
            full_response = "Ollama is offline. Start Ollama, then try again."
            yield full_response
        except requests.Timeout:
            full_response = (
                f"`{selected_model}` took too long to answer. "
                "Try Phi-3 Mini for faster local responses, or ask a shorter question."
            )
            yield full_response
        except requests.RequestException:
            full_response = f"I could not reach `{selected_model}` through Ollama. Choose an Active model and try again."
            yield full_response

        db.add_all([
            Message(id=str(uuid4()), chat_id=data.chat_id, role="user", content=data.message),
            Message(id=str(uuid4()), chat_id=data.chat_id, role="assistant", content=full_response),
        ])
        db.commit()

    return StreamingResponse(generate(), media_type="text/plain")
